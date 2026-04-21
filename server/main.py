import csv
import os
import re
from io import BytesIO

import pandas as pd
from docx import Document
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from openai import OpenAI
from openpyxl import load_workbook
from pydantic import BaseModel
from pptx import Presentation
from PyPDF2 import PdfReader

load_dotenv()

app = FastAPI(title="Mwakenya API")

origins = [
    "http://localhost:3000",
    "http://127.0.0.1:3000",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

api_key = os.getenv("OPENAI_API_KEY")
if not api_key:
    raise RuntimeError("OPENAI_API_KEY is missing. Add it to server/.env")

client = OpenAI(api_key=api_key)


class TextRequest(BaseModel):
    content: str


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def extract_pdf_text(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages).strip()


def extract_docx_text(file_bytes: bytes) -> str:
    document = Document(BytesIO(file_bytes))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def extract_pptx_text(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines).strip()


def extract_xlsx_text(file_bytes: bytes) -> str:
    workbook = load_workbook(BytesIO(file_bytes), data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines).strip()


def extract_txt_text(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore").strip()


def extract_csv_text(file_bytes: bytes) -> str:
    decoded = file_bytes.decode("utf-8", errors="ignore").splitlines()
    reader = csv.reader(decoded)
    rows = []
    for row in reader:
        if row:
            rows.append(" | ".join(row))
    return "\n".join(rows).strip()


def extract_text_from_upload(filename: str, file_bytes: bytes) -> str:
    lower_name = filename.lower()

    if lower_name.endswith(".pdf"):
        return extract_pdf_text(file_bytes)
    if lower_name.endswith(".docx"):
        return extract_docx_text(file_bytes)
    if lower_name.endswith(".pptx"):
        return extract_pptx_text(file_bytes)
    if lower_name.endswith(".xlsx"):
        return extract_xlsx_text(file_bytes)
    if lower_name.endswith(".txt"):
        return extract_txt_text(file_bytes)
    if lower_name.endswith(".csv"):
        return extract_csv_text(file_bytes)

    raise HTTPException(
        status_code=400,
        detail="Unsupported file type. Use PDF, DOCX, PPTX, XLSX, TXT, or CSV.",
    )


def truncate_content(text: str, max_chars: int = 12000) -> str:
    text = clean_text(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars]


@app.get("/")
def read_root():
    return {"message": "Mwakenya backend is running"}


@app.get("/health")
def health_check():
    return {"status": "ok"}


@app.post("/extract-text")
async def extract_text(file: UploadFile = File(...)):
    file_bytes = await file.read()

    if not file.filename:
        raise HTTPException(status_code=400, detail="No file uploaded.")

    extracted_text = extract_text_from_upload(file.filename, file_bytes)
    extracted_text = clean_text(extracted_text)

    if not extracted_text:
        raise HTTPException(
            status_code=400,
            detail="No readable text was found in the uploaded file.",
        )

    return {
        "filename": file.filename,
        "content": extracted_text,
    }


@app.post("/summarize")
def summarize_text(request: TextRequest):
    text = request.content.strip()

    if not text:
        raise HTTPException(status_code=400, detail="Content is required.")

    prepared_text = truncate_content(text)

    try:
        response = client.responses.create(
            model="gpt-4.1-mini",
            input=[
                {
                    "role": "system",
                    "content": (
                        "You are an expert educational AI tutor inside the "
                        "Adaptive Learning Intelligence Platform (ALIP). "
                        "Your job is to produce a strong student-friendly "
                        "summary for revision."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        "Read the study material below and produce a summary.\n\n"
                        "Rules:\n"
                        "1. Rewrite the material in your own words.\n"
                        "2. Make the summary clear, accurate, and useful for revision.\n"
                        "3. Focus on the main ideas, definitions, relationships, and applications.\n"
                        "4. Do not copy the original text unless necessary.\n"
                        "5. Do not invent facts that are not in the source.\n"
                        "6. Keep the summary detailed enough to be useful, but not too long.\n\n"
                        f"Study material:\n{prepared_text}"
                    ),
                },
            ],
        )

        summary = (response.output_text or "").strip()

        if not summary:
            raise HTTPException(status_code=500, detail="No summary returned.")

        return {"summary": summary}

    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Summary generation failed: {str(e)}",
        )